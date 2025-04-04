import logging
from flask import Flask, request, jsonify, render_template, send_from_directory, send_file
from flask_cors import CORS
from pptx import Presentation, util
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
import os
from huggingface_hub import InferenceClient
import random

app = Flask(__name__)
CORS(app)

# Configure logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(filename)s:%(lineno)d - %(message)s')

load_dotenv()

# Correctly retrieve the API token from environment variables
HUGGINGFACE_API_KEY = os.getenv("HF_API_TOKEN")  # Use os.getenv()
MODEL_ID = "mistralai/Mistral-7B-Instruct-v0.2"

# Check if the API key is loaded correctly
if HUGGINGFACE_API_KEY is None:
    logging.error("HF_API_TOKEN not found in environment variables. Please check your .env file.")
    exit(1)  # Exit if the API key is missing

client_hf = InferenceClient(token=HUGGINGFACE_API_KEY) # Use the variable that contains the API key

# --- Color Palettes ---
COLOR_PALETTES = [
    {  # Blue Theme
        "background": RGBColor(240, 248, 255),  # Light Alice Blue
        "text": RGBColor(0, 0, 128),  # Navy
        "accent": RGBColor(0, 191, 255),  # Deep Sky Blue
    },
    {  # Green Theme
        "background": RGBColor(240, 255, 240),  # Honeydew
        "text": RGBColor(0, 100, 0),  # Dark Green
        "accent": RGBColor(50, 205, 50),  # Lime Green
    },
    {  # Purple Theme
        "background": RGBColor(240, 240, 255),  # Light Lavender
        "text": RGBColor(75, 0, 130),  # Indigo
        "accent": RGBColor(148, 0, 211),  # Dark Violet
    }
]

# --- Font Families ---
TITLE_FONT = "Arial Black"
BODY_FONT = "Calibri"

# --- Layout Functions ---
def apply_centered_layout(slide, title, content, color_palette):
    """Applies a centered layout to the slide."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_palette["background"]
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = color_palette["text"]

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = BODY_FONT
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = color_palette["text"]
    tf.text = content

def apply_bullet_point_layout(slide, title, bullet_points, color_palette):
    """Applies a bullet point layout to the slide."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_palette["background"]
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = color_palette["text"]

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].font.name = BODY_FONT
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = color_palette["text"]

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

def apply_three_point_circle_layout(slide, title, points, color_palette):
    """Applies a three-point circle layout to the slide."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_palette["background"]
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = color_palette["text"]

    # Define circle positions
    circle_positions = [
        (Inches(4), Inches(1)),  # Top
        (Inches(1), Inches(4)),  # Bottom Left
        (Inches(7), Inches(4)),  # Bottom Right
    ]

    for i, point in enumerate(points):
        left, top = circle_positions[i]
        width = Inches(2)
        height = Inches(2)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_palette["accent"]
        shape.line.fill.background()

        text_frame = shape.text_frame
        text_frame.text = point
        text_frame.word_wrap = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.name = BODY_FONT
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.color.rgb = color_palette["text"]
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def apply_two_column_layout(slide, title, content1, content2, color_palette):
    """Applies a two-column layout to the slide."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_palette["background"]
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = color_palette["text"]

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(5)

    textbox1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = textbox1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].font.name = BODY_FONT
    tf1.paragraphs[0].font.size = Pt(20)
    tf1.paragraphs[0].font.color.rgb = color_palette["text"]
    tf1.text = content1

    textbox2 = slide.shapes.add_textbox(left + Inches(5), top, width, height)
    tf2 = textbox2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].font.name = BODY_FONT
    tf2.paragraphs[0].font.size = Pt(20)
    tf2.paragraphs[0].font.color.rgb = color_palette["text"]
    tf2.text = content2

def apply_four_point_layout(slide, title, points, color_palette):
    """Applies a four-point grid layout to the slide."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_palette["background"]
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = color_palette["text"]

    # Define grid positions
    grid_positions = [
        (Inches(1), Inches(1.5)),  # Top Left
        (Inches(6), Inches(1.5)),  # Top Right
        (Inches(1), Inches(4)),  # Bottom Left
        (Inches(6), Inches(4)),  # Bottom Right
    ]

    for i, point in enumerate(points):
        left, top = grid_positions[i]
        width = Inches(4)
        height = Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].font.name = BODY_FONT
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = color_palette["text"]
        tf.text = point

def apply_text_only_layout(slide, title, content, color_palette):
    """Applies a text-only layout to the slide."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color_palette["background"]
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = color_palette["text"]

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].font.name = BODY_FONT
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = color_palette["text"]
    tf.text = content

def generate_presentation_outline(prompt, num_slides=10):
    """Generates a presentation outline using a Hugging Face model."""
    logging.info(f"Generating presentation outline for prompt: '{prompt}'")
    try:
        response = client_hf.text_generation(
            prompt=f"""Create a detailed presentation outline on {prompt}.  The outline should include a title, and exactly {num_slides} slides with clear titles. Each slide should have 0-5 bullet points. Use concise and informative bullet points. When it is relevant, add examples to the bullet points. Structure the output as follows:

Title: [Presentation Title]

Slide 1: [Slide Title 1]
- Bullet Point 1
- Bullet Point 2
- Bullet Point 3
- Example: [Example]

Slide 2: [Slide Title 2]
- Bullet Point 1
- Bullet Point 2

Slide 3: [Slide Title 3]
- Brief discussion of the topic

... (and so on for {num_slides} slides)
""",
            model=MODEL_ID,
            max_new_tokens=1000, # Increased tokens for more slides
        )
        logging.debug(f"Hugging Face API response: {response}")
        generated_content = response

        lines = generated_content.splitlines()
        title = ""
        slides = []
        current_slide = None
        slide_count = 0

        for line in lines:
            line = line.strip()
            if line.startswith("Title:"):
                title = line[len("Title:"):].strip()
            elif line.startswith("Slide"):
                if current_slide:
                    slides.append(current_slide)
                try:
                    current_slide = {"title": line.split(":")[1].strip(), "bullet_points": []}
                    slide_count +=1
                    if slide_count > num_slides:
                        break
                except IndexError:
                    logging.error(f"Error parsing slide line: {line}")
                    return None
            elif line.startswith("-") and current_slide:
                current_slide["bullet_points"].append(line[1:].strip())
            elif line and current_slide:
                current_slide["bullet_points"].append(line)

        if current_slide:
            slides.append(current_slide)
        
        # Ensure we have exactly 10 slides, add empty ones if needed
        while len(slides) < num_slides:
            slides.append({"title": "Slide " + str(len(slides) + 1), "bullet_points": []})

        return {"title": title, "slides": slides}

    except Exception as e:
        logging.exception(f"Unexpected error in generate_presentation_outline: {e}")
        return None

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    logging.info("Received POST request to /generate_ppt")
    try:
        data = request.get_json()
        prompt = data.get('prompt')
        logging.debug(f"Received prompt: {prompt}")

        if not prompt:
            return jsonify({'error': 'No prompt provided.'}), 400

        outline = generate_presentation_outline(prompt, num_slides=10) # Generate 10 slides

        if not outline:
            return jsonify({'error': 'Failed to generate presentation outline.'}), 500

        title = outline["title"]
        slides = outline["slides"]

        prs = Presentation()
        # --- Apply a Random Theme ---
        color_palette = random.choice(COLOR_PALETTES)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # --- Title Slide ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color_palette["background"]
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = TITLE_FONT
        title_shape.text_frame.paragraphs[0].font.size = Pt(48)
        title_shape.text_frame.paragraphs[0].font.color.rgb = color_palette["text"]

        # --- Content Slides ---
        for slide_data in slides:
            slide_title = slide_data.get('title', '')
            bullet_points = slide_data.get('bullet_points', [])

            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # --- Choose Layout Based on Content ---
            num_points = len(bullet_points)
            layout_choice = random.randint(0, 100)

            if num_points == 0:
                apply_text_only_layout(slide, slide_title, " ".join(bullet_points), color_palette)
            elif num_points == 1:
                apply_centered_layout(slide, slide_title, bullet_points[0], color_palette)
            elif num_points == 2:
                apply_two_column_layout(slide, slide_title, bullet_points[0], bullet_points[1], color_palette)
            elif num_points == 3:
                if layout_choice < 70:
                    apply_three_point_circle_layout(slide, slide_title, bullet_points, color_palette)
                else:
                    apply_bullet_point_layout(slide, slide_title, bullet_points, color_palette)
            elif num_points == 4:
                if layout_choice < 70:
                    apply_four_point_layout(slide, slide_title, bullet_points, color_palette)
                else:
                    apply_bullet_point_layout(slide, slide_title, bullet_points, color_palette)
            else:
                apply_bullet_point_layout(slide, slide_title, bullet_points, color_palette)

        prs.save('presentation.pptx')
        return jsonify({'message': 'Presentation generated successfully!'})

    except Exception as e:
        logging.exception(f"Unexpected error in /generate_ppt: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>')
def download_file(filename):
    try:
        return send_file(filename, as_attachment=True, download_name=filename)
    except FileNotFoundError:
        return jsonify({'error': 'File not found.'}), 404

@app.route('/test_post', methods=['POST'])
def test_post():
    logging.info("Received POST request to /test_post")
    return jsonify({'message': 'POST request received!'}), 200

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')

@app.route('/create_presentation.html')
def create_presentation():
    return send_from_directory('.', 'create_presentation.html')

@app.route('/<path:filename>')
def serve_static(filename):
    return send_from_directory('.', filename)

if __name__ == '__main__':
    app.run(debug=True, port=8000)
